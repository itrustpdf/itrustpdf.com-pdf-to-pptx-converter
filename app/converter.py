"""
Main PDF to PPTX and PPTX to PDF conversion pipeline.
"""

import fitz  # PyMuPDF
from typing import List, Optional
import logging
import io
import tempfile
import os
import subprocess
import glob

from pptx import Presentation
from pptx.util import Inches
from reportlab.pdfgen import canvas
from PIL import Image, ImageDraw, ImageFont
import textwrap

from .models import TextBlock, MINIMUM_TEXT_THRESHOLD
from .utils import get_pdf_dimensions
from .text_extraction import extract_text_blocks_pymupdf, has_sufficient_text, normalize_and_group_text_blocks
from .ocr import ocr_page_lines
from .layout import transform_blocks_to_pptx
from .pptx_generator import create_pptx_from_blocks, calculate_optimal_slide_size

logger = logging.getLogger(__name__)


def pdf_to_pptx(pdf_bytes: bytes, 
               ocr_langs: str = 'eng', 
               dehyphenate: bool = True,
               use_ocr: bool = True) -> bytes:
    """
    Convert PDF bytes to PPTX bytes.
    
    Args:
        pdf_bytes: PDF file content as bytes
        ocr_langs: Tesseract language codes for OCR
        dehyphenate: Whether to remove end-of-line hyphenation
        use_ocr: If True, extract text with OCR. If False, convert pages to images.
        
    Returns:
        PPTX file content as bytes
        
    Raises:
        ValueError: If PDF processing fails
        Exception: If conversion fails
    """
    if use_ocr:
        return _pdf_to_pptx_with_ocr(pdf_bytes, ocr_langs, dehyphenate)
    else:
        return _pdf_to_pptx_as_images(pdf_bytes)


def _pdf_to_pptx_with_ocr(pdf_bytes: bytes, 
                         ocr_langs: str = 'eng', 
                         dehyphenate: bool = True) -> bytes:
    """
    Convert PDF to PPTX using OCR to extract and preserve text formatting.
    
    Args:
        pdf_bytes: PDF file content as bytes
        ocr_langs: Tesseract language codes for OCR
        dehyphenate: Whether to remove end-of-line hyphenation
        
    Returns:
        PPTX file content as bytes
    """
    try:
        logger.info("Starting PDF to PPTX conversion with OCR")
        
        # Open PDF document
        doc = fitz.open(stream=pdf_bytes, filetype="pdf")
        
        if len(doc) == 0:
            raise ValueError("Empty PDF")
        
        # Get PDF dimensions from first page
        first_page = doc[0]
        pdf_width = first_page.rect.width
        pdf_height = first_page.rect.height
        page_count = len(doc)
        
        logger.info(f"Processing PDF: {page_count} pages, {pdf_width}x{pdf_height} points")
        
        # Calculate optimal slide configuration
        slide_config = calculate_optimal_slide_size(pdf_width, pdf_height)
        
        # Process each page
        all_page_blocks = []
        
        for page_num in range(page_count):
            page = doc[page_num]
            logger.info(f"Processing page {page_num + 1}/{page_count}")
            
            # Extract text blocks for this page
            page_blocks = _extract_page_text_blocks(page, ocr_langs)
            
            # Normalize and group text blocks
            normalized_blocks = normalize_and_group_text_blocks(page_blocks, dehyphenate)
            
            # Transform to PPTX coordinates
            transformed_blocks = transform_blocks_to_pptx(
                normalized_blocks, pdf_width, pdf_height, slide_config
            )
            
            all_page_blocks.append(transformed_blocks)
            logger.info(f"Page {page_num + 1}: {len(transformed_blocks)} text blocks")
        
        doc.close()
        
        # Generate PPTX
        pptx_bytes = create_pptx_from_blocks(all_page_blocks, slide_config)
        
        logger.info(f"OCR conversion completed: {len(pptx_bytes)} bytes")
        return pptx_bytes
        
    except Exception as e:
        logger.error(f"PDF to PPTX with OCR failed: {str(e)}")
        raise Exception(f"OCR conversion failed: {str(e)}")


def _pdf_to_pptx_as_images(pdf_bytes: bytes) -> bytes:
    """
    Convert PDF to PPTX by placing each page as an image on a slide.
    
    Args:
        pdf_bytes: PDF file content as bytes
        
    Returns:
        PPTX file content as bytes
    """
    try:
        logger.info("Starting PDF to PPTX conversion as images")
        
        # Create temporary directory
        temp_dir = tempfile.mkdtemp()
        
        try:
            # Save PDF to temporary file
            pdf_path = os.path.join(temp_dir, "input.pdf")
            with open(pdf_path, 'wb') as f:
                f.write(pdf_bytes)
            
            # Convert PDF pages to images
            image_paths = _convert_pdf_to_images(pdf_path, temp_dir)
            
            if not image_paths:
                raise ValueError("Failed to convert PDF pages to images")
            
            logger.info(f"Converted {len(image_paths)} pages to images")
            
            # Create a new presentation
            from pptx import Presentation
            from pptx.util import Inches
            
            presentation = Presentation()
            
            # Get slide dimensions from first image
            with Image.open(image_paths[0]) as img:
                img_width, img_height = img.size
            
            # Calculate aspect ratio
            aspect_ratio = img_width / img_height
            
            # Set slide size based on image aspect ratio
            # Standard slide size is 10x7.5 inches (4:3) or 13.33x7.5 inches (16:9)
            if aspect_ratio > 1.5:  # Wider than 3:2, use 16:9
                slide_width = Inches(13.33)
                slide_height = Inches(7.5)
            else:  # Use 4:3
                slide_width = Inches(10)
                slide_height = Inches(7.5)
            
            presentation.slide_width = slide_width
            presentation.slide_height = slide_height
            
            # Add each image as a slide
            for i, image_path in enumerate(sorted(image_paths)):
                logger.info(f"Adding page {i + 1}/{len(image_paths)} as slide")
                
                # Add a blank slide
                slide_layout = presentation.slide_layouts[6]  # Blank layout
                slide = presentation.slides.add_slide(slide_layout)
                
                # Add image to slide
                left = Inches(0.5)
                top = Inches(0.5)
                width = slide_width - Inches(1)  # 1 inch margins
                height = slide_height - Inches(1)
                
                pic = slide.shapes.add_picture(image_path, left, top, width, height)
            
            # Save presentation to bytes
            pptx_buffer = io.BytesIO()
            presentation.save(pptx_buffer)
            pptx_bytes = pptx_buffer.getvalue()
            
            logger.info(f"Image conversion completed: {len(pptx_bytes)} bytes")
            return pptx_bytes
            
        finally:
            # Clean up
            import shutil
            try:
                shutil.rmtree(temp_dir)
            except:
                pass
            
    except Exception as e:
        logger.error(f"PDF to PPTX as images failed: {str(e)}")
        raise Exception(f"Image conversion failed: {str(e)}")


def _convert_pdf_to_images(pdf_path: str, output_dir: str) -> List[str]:
    """
    Convert PDF pages to images using PyMuPDF.
    
    Args:
        pdf_path: Path to PDF file
        output_dir: Directory to save images
        
    Returns:
        List of paths to generated image files
    """
    try:
        image_paths = []
        doc = fitz.open(pdf_path)
        
        for page_num in range(len(doc)):
            page = doc[page_num]
            
            # Get page dimensions
            rect = page.rect
            zoom = 2.0  # Zoom factor for better quality
            mat = fitz.Matrix(zoom, zoom)
            
            # Render page to image
            pix = page.get_pixmap(matrix=mat, alpha=False)
            
            # Save image
            image_path = os.path.join(output_dir, f"page_{page_num + 1:03d}.png")
            pix.save(image_path)
            
            image_paths.append(image_path)
            logger.info(f"Saved page {page_num + 1} as image: {image_path}")
        
        doc.close()
        return image_paths
        
    except Exception as e:
        logger.error(f"Failed to convert PDF to images: {str(e)}")
        
        # Try alternative method using pdftoppm if available
        try:
            return _convert_pdf_to_images_pdftoppm(pdf_path, output_dir)
        except Exception as e2:
            logger.error(f"pdftoppm also failed: {str(e2)}")
            return []


def _convert_pdf_to_images_pdftoppm(pdf_path: str, output_dir: str) -> List[str]:
    """
    Convert PDF pages to images using pdftoppm (fallback).
    
    Args:
        pdf_path: Path to PDF file
        output_dir: Directory to save images
        
    Returns:
        List of paths to generated image files
    """
    try:
        # Check if pdftoppm is available
        result = subprocess.run(['which', 'pdftoppm'], 
                              capture_output=True, text=True)
        if result.returncode != 0:
            logger.warning("pdftoppm not found in PATH")
            return []
        
        # Convert PDF to PNG using pdftoppm
        output_pattern = os.path.join(output_dir, "page")
        cmd = [
            'pdftoppm',
            '-png',
            '-r', '150',  # 150 DPI for good quality
            pdf_path,
            output_pattern
        ]
        
        logger.info(f"Running pdftoppm: {' '.join(cmd)}")
        result = subprocess.run(cmd, capture_output=True, text=True, timeout=60)
        
        if result.returncode == 0:
            # Find generated PNG files
            image_paths = glob.glob(os.path.join(output_dir, "page*.png"))
            image_paths.sort()
            logger.info(f"pdftoppm converted {len(image_paths)} pages")
            return image_paths
        else:
            logger.warning(f"pdftoppm failed: {result.stderr}")
            return []
            
    except Exception as e:
        logger.warning(f"pdftoppm conversion error: {str(e)}")
        return []


def pptx_to_pdf(pptx_bytes: bytes) -> bytes:
    """
    Convert PPTX bytes to PDF bytes by converting each slide to an image.
    
    Args:
        pptx_bytes: PPTX file content as bytes
        
    Returns:
        PDF file content as bytes
        
    Raises:
        ValueError: If PPTX processing fails
        Exception: If conversion fails
    """
    try:
        logger.info("Starting PPTX to PDF conversion")
        
        # Create temporary directory for all files
        temp_dir = tempfile.mkdtemp()
        
        try:
            # Save PPTX to temporary file
            pptx_path = os.path.join(temp_dir, "input.pptx")
            with open(pptx_path, 'wb') as f:
                f.write(pptx_bytes)
            
            # Convert PPTX to images using LibreOffice
            image_paths = _convert_pptx_to_images_libreoffice(pptx_path, temp_dir)
            
            if not image_paths:
                # Fallback: Use python-pptx method
                logger.info("LibreOffice conversion failed, using fallback method")
                image_paths = _convert_pptx_to_images_fallback(pptx_path, temp_dir)
            
            if not image_paths:
                raise ValueError("Failed to convert any slides to images")
            
            logger.info(f"Found {len(image_paths)} images to convert to PDF")
            
            # Get dimensions from first image
            with Image.open(image_paths[0]) as img:
                img_width, img_height = img.size
            
            # Standard PDF DPI is 72, images are typically 96 DPI
            # Convert image pixels to PDF points
            pdf_width = img_width * (72 / 96)
            pdf_height = img_height * (72 / 96)
            
            # Create PDF from images
            pdf_buffer = io.BytesIO()
            c = canvas.Canvas(pdf_buffer, pagesize=(pdf_width, pdf_height))
            
            for i, image_path in enumerate(sorted(image_paths)):
                logger.info(f"Adding slide {i + 1}/{len(image_paths)} to PDF")
                
                # Add image to PDF page
                c.drawImage(image_path, 0, 0, pdf_width, pdf_height)
                
                # Add new page for next slide (except last one)
                if i < len(image_paths) - 1:
                    c.showPage()
            
            # Save PDF
            c.save()
            pdf_bytes = pdf_buffer.getvalue()
            
            logger.info(f"Conversion completed: {len(pdf_bytes)} bytes")
            return pdf_bytes
            
        finally:
            # Clean up temporary directory
            import shutil
            try:
                shutil.rmtree(temp_dir)
            except:
                pass
            
    except Exception as e:
        logger.error(f"PPTX to PDF conversion failed: {str(e)}")
        raise Exception(f"Conversion failed: {str(e)}")


def _convert_pptx_to_images_libreoffice(pptx_path: str, output_dir: str) -> List[str]:
    """
    Convert PPTX slides to images using LibreOffice.
    
    Args:
        pptx_path: Path to PPTX file
        output_dir: Directory to save images
        
    Returns:
        List of paths to generated image files
    """
    try:
        # Check if LibreOffice is available
        result = subprocess.run(['which', 'libreoffice'], 
                              capture_output=True, text=True)
        if result.returncode != 0:
            logger.warning("LibreOffice not found in PATH")
            return []
        
        # Convert PPTX to PNG using LibreOffice
        cmd = [
            'libreoffice',
            '--headless',
            '--convert-to', 'png',
            '--outdir', output_dir,
            pptx_path
        ]
        
        logger.info(f"Running LibreOffice conversion: {' '.join(cmd)}")
        result = subprocess.run(cmd, capture_output=True, text=True, timeout=60)
        
        if result.returncode == 0:
            logger.info("LibreOffice conversion successful")
            
            # Find all PNG files in output directory
            image_paths = glob.glob(os.path.join(output_dir, "*.png"))
            
            # Sort by filename to maintain slide order
            image_paths.sort()
            
            logger.info(f"Found {len(image_paths)} PNG files")
            return image_paths
        else:
            logger.warning(f"LibreOffice conversion failed: {result.stderr}")
            return []
            
    except subprocess.TimeoutExpired:
        logger.warning("LibreOffice conversion timed out")
        return []
    except Exception as e:
        logger.warning(f"LibreOffice conversion error: {str(e)}")
        return []


def _convert_pptx_to_images_fallback(pptx_path: str, output_dir: str) -> List[str]:
    """
    Fallback method to convert PPTX to images using python-pptx.
    
    Args:
        pptx_path: Path to PPTX file
        output_dir: Directory to save images
        
    Returns:
        List of paths to generated image files
    """
    try:
        presentation = Presentation(pptx_path)
        image_paths = []
        
        for slide_idx, slide in enumerate(presentation.slides):
            # Create image for this slide
            image_path = os.path.join(output_dir, f"slide_{slide_idx + 1:03d}.png")
            
            # Get slide dimensions
            slide_width = presentation.slide_width.inches
            slide_height = presentation.slide_height.inches
            
            # Convert to pixels (96 DPI)
            width_px = int(slide_width * 96)
            height_px = int(slide_height * 96)
            
            # Create slide image
            _create_slide_image(slide, slide_idx + 1, width_px, height_px, image_path)
            
            image_paths.append(image_path)
            logger.info(f"Created fallback image for slide {slide_idx + 1}")
        
        return image_paths
        
    except Exception as e:
        logger.error(f"Fallback conversion failed: {str(e)}")
        return []


def _create_slide_image(slide, slide_num: int, width_px: int, height_px: int, output_path: str):
    """
    Create an image representation of a slide.
    
    Args:
        slide: PPTX slide object
        slide_num: Slide number (1-based)
        width_px: Image width in pixels
        height_px: Image height in pixels
        output_path: Path to save the image
    """
    try:
        # Create a white background
        img = Image.new('RGB', (width_px, height_px), color='white')
        draw = ImageDraw.Draw(img)
        
        # Try to load fonts
        try:
            font_large = ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf", 28)
            font_medium = ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf", 18)
            font_small = ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf", 14)
        except:
            # Fallback to default font
            font_large = ImageFont.load_default()
            font_medium = ImageFont.load_default()
            font_small = ImageFont.load_default()
        
        # Draw slide header
        title = f"Slide {slide_num}"
        title_width = draw.textlength(title, font=font_large)
        draw.text(
            ((width_px - title_width) // 2, 50),
            title,
            fill='darkblue',
            font=font_large
        )
        
        # Draw separator line
        draw.line([(50, 100), (width_px - 50, 100)], fill='gray', width=2)
        
        # Extract and draw text from shapes
        y_offset = 130
        shapes_with_text = []
        
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                shapes_with_text.append(shape)
        
        # Limit to avoid overflow
        max_shapes = min(10, len(shapes_with_text))
        
        for i in range(max_shapes):
            shape = shapes_with_text[i]
            text = shape.text.strip()
            
            if not text:
                continue
            
            # Truncate long text
            if len(text) > 100:
                text = text[:97] + "..."
            
            # Wrap text
            wrapped_lines = textwrap.wrap(text, width=50)
            
            for line in wrapped_lines:
                if y_offset < height_px - 50:
                    line_width = draw.textlength(line, font=font_small)
                    draw.text(
                        ((width_px - line_width) // 2, y_offset),
                        line,
                        fill='black',
                        font=font_small
                    )
                    y_offset += 25
        
        # If no text was found, add a message
        if not shapes_with_text:
            message = "Slide contains no extractable text"
            msg_width = draw.textlength(message, font=font_medium)
            draw.text(
                ((width_px - msg_width) // 2, height_px // 2),
                message,
                fill='gray',
                font=font_medium
            )
        
        # Save the image
        img.save(output_path, 'PNG', quality=90)
        
    except Exception as e:
        logger.error(f"Failed to create slide image: {str(e)}")
        # Create a simple fallback image
        img = Image.new('RGB', (width_px, height_px), color='lightgray')
        draw = ImageDraw.Draw(img)
        
        try:
            font = ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf", 24)
        except:
            font = ImageFont.load_default()
        
        text = f"Slide {slide_num}"
        text_width = draw.textlength(text, font=font)
        draw.text(
            ((width_px - text_width) // 2, height_px // 2 - 20),
            text,
            fill='darkred',
            font=font
        )
        
        msg = "Content could not be rendered"
        msg_width = draw.textlength(msg, font=font)
        draw.text(
            ((width_px - msg_width) // 2, height_px // 2 + 20),
            msg,
            fill='darkred',
            font=font
        )
        
        img.save(output_path, 'PNG')


def _extract_page_text_blocks(page: fitz.Page, ocr_langs: str) -> List[TextBlock]:
    """
    Extract text blocks from a single PDF page using native extraction or OCR.
    
    Args:
        page: PyMuPDF page object
        ocr_langs: Language codes for OCR
        
    Returns:
        List of text blocks
    """
    try:
        # First, try native text extraction
        native_blocks = extract_text_blocks_pymupdf(page)
        
        # Check if we have sufficient text
        if has_sufficient_text(native_blocks):
            logger.debug(f"Using native text extraction: {len(native_blocks)} blocks")
            return native_blocks
        else:
            logger.debug(f"Insufficient native text ({sum(len(b[4]) for b in native_blocks)} chars), using OCR")
            
        # Fall back to OCR
        try:
            ocr_blocks = ocr_page_lines(page, langs=ocr_langs)
            logger.debug(f"OCR extracted {len(ocr_blocks)} blocks")
            return ocr_blocks
            
        except Exception as ocr_error:
            logger.warning(f"OCR failed: {str(ocr_error)}, using native extraction")
            return native_blocks
            
    except Exception as e:
        logger.error(f"Text extraction failed for page: {str(e)}")
        return []  # Return empty blocks for failed pages


def validate_pdf(pdf_bytes: bytes) -> bool:
    """
    Validate that the input is a valid PDF.
    
    Args:
        pdf_bytes: PDF file content as bytes
        
    Returns:
        True if valid PDF, False otherwise
    """
    try:
        doc = fitz.open(stream=pdf_bytes, filetype="pdf")
        is_valid = len(doc) > 0
        doc.close()
        return is_valid
    except Exception as e:
        logger.error(f"PDF validation failed: {str(e)}")
        return False


def validate_pptx(pptx_bytes: bytes) -> bool:
    """
    Validate that the input is a valid PPTX file.
    
    Args:
        pptx_bytes: PPTX file content as bytes
        
    Returns:
        True if valid PPTX, False otherwise
    """
    try:
        with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp_file:
            tmp_file.write(pptx_bytes)
            tmp_path = tmp_file.name
        
        try:
            presentation = Presentation(tmp_path)
            # Check if we can access basic properties
            _ = len(presentation.slides)
            return True
        except Exception as e:
            logger.error(f"PPTX validation failed: {str(e)}")
            return False
        finally:
            try:
                os.unlink(tmp_path)
            except:
                pass
            
    except Exception as e:
        logger.error(f"PPTX file handling failed: {str(e)}")
        return False


def get_pdf_info(pdf_bytes: bytes) -> dict:
    """
    Extract basic information from a PDF.
    
    Args:
        pdf_bytes: PDF file content as bytes
        
    Returns:
        Dictionary with PDF information
    """
    try:
        doc = fitz.open(stream=pdf_bytes, filetype="pdf")
        
        metadata = doc.metadata or {}
        info = {
            'page_count': len(doc),
            'title': metadata.get('title', ''),
            'author': metadata.get('author', ''),
            'subject': metadata.get('subject', ''),
            'creator': metadata.get('creator', ''),
            'producer': metadata.get('producer', ''),
            'creation_date': metadata.get('creationDate', ''),
            'modification_date': metadata.get('modDate', ''),
        }
        
        if len(doc) > 0:
            first_page = doc[0]
            info['page_width'] = first_page.rect.width
            info['page_height'] = first_page.rect.height
            info['page_aspect_ratio'] = first_page.rect.width / first_page.rect.height
        
        doc.close()
        return info
        
    except Exception as e:
        logger.error(f"Failed to get PDF info: {str(e)}")
        return {'error': str(e)}


def get_pptx_info(pptx_bytes: bytes) -> dict:
    """
    Extract basic information from a PPTX file.
    
    Args:
        pptx_bytes: PPTX file content as bytes
        
    Returns:
        Dictionary with PPTX information
    """
    try:
        with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp_file:
            tmp_file.write(pptx_bytes)
            tmp_path = tmp_file.name
        
        try:
            presentation = Presentation(tmp_path)
            
            info = {
                'slide_count': len(presentation.slides),
                'slide_width_inches': presentation.slide_width.inches,
                'slide_height_inches': presentation.slide_height.inches,
                'slide_width_points': presentation.slide_width.inches * 72,
                'slide_height_points': presentation.slide_height.inches * 72,
                'slide_aspect_ratio': presentation.slide_width.inches / presentation.slide_height.inches,
            }
            
            return info
            
        finally:
            try:
                os.unlink(tmp_path)
            except:
                pass
            
    except Exception as e:
        logger.error(f"Failed to get PPTX info: {str(e)}")
        return {'error': str(e)}


def estimate_processing_time(pdf_bytes: bytes, use_ocr: bool = True) -> float:
    """
    Estimate processing time for a PDF based on page count and content complexity.
    
    Args:
        pdf_bytes: PDF file content as bytes
        use_ocr: Whether OCR will be used
        
    Returns:
        Estimated processing time in seconds
    """
    try:
        doc = fitz.open(stream=pdf_bytes, filetype="pdf")
        page_count = len(doc)
        doc.close()
        
        if use_ocr:
            # OCR mode
            base_time_per_page = 2.0
            ocr_time_per_page = 5.0
            # Assume 50% of pages might need OCR
            estimated_time = (page_count * base_time_per_page) + (page_count * 0.5 * ocr_time_per_page)
        else:
            # Image mode (faster)
            estimated_time = page_count * 1.0
        
        return max(3.0, estimated_time)  # Minimum 3 seconds
        
    except Exception as e:
        logger.error(f"Failed to estimate processing time: {str(e)}")
        return 30.0  # Default estimate


def estimate_pptx_processing_time(pptx_bytes: bytes) -> float:
    """
    Estimate processing time for a PPTX based on slide count.
    
    Args:
        pptx_bytes: PPTX file content as bytes
        
    Returns:
        Estimated processing time in seconds
    """
    try:
        info = get_pptx_info(pptx_bytes)
        slide_count = info.get('slide_count', 1)
        
        # Base time per slide (seconds)
        base_time_per_slide = 3.0
        
        estimated_time = slide_count * base_time_per_slide
        
        return max(3.0, estimated_time)  # Minimum 3 seconds
        
    except Exception as e:
        logger.error(f"Failed to estimate PPTX processing time: {str(e)}")
        return 15.0  # Default estimate
