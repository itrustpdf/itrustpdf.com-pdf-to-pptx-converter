"""
PPTX generation module for creating natural PowerPoint presentations from text content.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.dml.color import RGBColor
from typing import List, Tuple
import io
import logging
import re

from .models import SlideConfig, WIDESCREEN_ASPECT_RATIO

logger = logging.getLogger(__name__)


def create_pptx_from_blocks(text_blocks_by_page: List[List[Tuple[int, int, int, int, str]]], 
                           slide_config: SlideConfig) -> bytes:
    """
    Create a natural PowerPoint presentation from text blocks organized by page.
    
    Args:
        text_blocks_by_page: List of text block lists, one per page
        slide_config: Slide configuration
        
    Returns:
        PPTX file as bytes
        
    Raises:
        Exception: If PPTX generation fails
    """
    try:
        # Create presentation with standard layout
        prs = Presentation()
        
        logger.info(f"Creating natural PowerPoint with {len(text_blocks_by_page)} slides")
        
        # Process each page into natural content
        for page_num, page_blocks in enumerate(text_blocks_by_page):
            _create_natural_slide(prs, page_blocks, page_num + 1)
        
        # Save to bytes
        pptx_bytes = io.BytesIO()
        prs.save(pptx_bytes)
        pptx_bytes.seek(0)
        
        logger.info("Natural PowerPoint generation completed successfully")
        return pptx_bytes.getvalue()
        
    except Exception as e:
        logger.error(f"PowerPoint generation failed: {str(e)}")
        raise Exception(f"PowerPoint generation failed: {str(e)}")


def _create_natural_slide(prs: Presentation, 
                         text_blocks: List[Tuple[int, int, int, int, str]], 
                         page_number: int) -> None:
    """
    Create a natural PowerPoint slide from text blocks.
    
    Args:
        prs: PowerPoint presentation object
        text_blocks: List of text blocks with coordinates
        page_number: Page number for logging
    """
    # Use title and content layout for natural presentation
    title_content_layout = prs.slide_layouts[1]  # Title and Content layout
    slide = prs.slides.add_slide(title_content_layout)
    
    if not text_blocks:
        logger.info(f"Created blank slide {page_number}")
        return
    
    # Extract and process text content naturally
    content_text = _extract_natural_content(text_blocks)
    
    if not content_text.strip():
        logger.info(f"No content found for slide {page_number}")
        return
    
    # Determine if there's a clear title
    title, body_content = _extract_title_and_content(content_text)
    
    # Set title if we found one
    if title and len(slide.shapes) > 0:
        title_shape = slide.shapes.title
        if title_shape:
            title_shape.text = title
    
    # Add body content
    if body_content.strip() and len(slide.shapes) > 1:
        content_placeholder = None
        for shape in slide.shapes:
            if hasattr(shape, 'text_frame') and shape != slide.shapes.title:
                content_placeholder = shape
                break
        
        if content_placeholder:
            _add_natural_content(content_placeholder, body_content)
        else:
            # Fallback: create manual text box
            _add_manual_content_box(slide, body_content)
    else:
        # No title found, treat all as content
        if len(slide.shapes) > 1:
            content_placeholder = None
            for shape in slide.shapes:
                if hasattr(shape, 'text_frame') and shape != slide.shapes.title:
                    content_placeholder = shape
                    break
            
            if content_placeholder:
                _add_natural_content(content_placeholder, content_text)
            else:
                _add_manual_content_box(slide, content_text)
    
    logger.info(f"Created natural slide {page_number}")


def _extract_natural_content(text_blocks: List[Tuple[int, int, int, int, str]]) -> str:
    """
    Extract text content in natural reading order.
    
    Args:
        text_blocks: List of text blocks with coordinates
        
    Returns:
        Combined text content in natural flow
    """
    if not text_blocks:
        return ""
    
    # Sort blocks by reading order (top to bottom, left to right)
    sorted_blocks = sorted(text_blocks, key=lambda b: (b[1], b[0]))
    
    content_parts = []
    for _, _, _, _, text in sorted_blocks:
        text = text.strip()
        if text:
            content_parts.append(text)
    
    # Combine with intelligent spacing
    return _intelligent_text_combination(content_parts)


def _intelligent_text_combination(text_parts: List[str]) -> str:
    """
    Combine text parts with intelligent spacing and formatting.
    
    Args:
        text_parts: List of text strings
        
    Returns:
        Intelligently combined text
    """
    if not text_parts:
        return ""
    
    if len(text_parts) == 1:
        return text_parts[0]
    
    result = []
    for i, part in enumerate(text_parts):
        if i == 0:
            result.append(part)
        else:
            prev_part = text_parts[i-1]
            
            # Determine spacing
            if _needs_paragraph_break(prev_part, part):
                result.append('\n\n' + part)
            elif _needs_line_break(prev_part, part):
                result.append('\n' + part)
            else:
                # Just add space if previous doesn't end with space
                if not prev_part.endswith(' '):
                    result.append(' ' + part)
                else:
                    result.append(part)
    
    return ''.join(result)


def _needs_paragraph_break(prev_text: str, current_text: str) -> bool:
    """Determine if paragraph break is needed."""
    # After sentence endings
    if prev_text.rstrip().endswith(('.', '!', '?', ':')):
        return True
    
    # Before bullet points or numbered items
    if re.match(r'^[\d•\-\*]\s*', current_text):
        return True
    
    # Between different formatting (all caps, etc.)
    if current_text.isupper() and len(current_text) > 10:
        return True
    
    # Before section headers
    if current_text.startswith(('Chapter', 'Section', 'Part', 'Article')):
        return True
    
    return False


def _needs_line_break(prev_text: str, current_text: str) -> bool:
    """Determine if line break is needed."""
    # Different formatting emphasis
    if current_text.isupper() or current_text.startswith('('):
        return True
    
    # After commas or semicolons in long text
    if len(prev_text) > 50 and prev_text.rstrip().endswith((',', ';')):
        return True
    
    return False


def _extract_title_and_content(text: str) -> Tuple[str, str]:
    """
    Extract potential title and content from text.
    
    Args:
        text: Full text content
        
    Returns:
        Tuple of (title, remaining_content)
    """
    lines = text.split('\n')
    if not lines:
        return "", ""
    
    first_line = lines[0].strip()
    
    # Check if first line looks like a title
    if _looks_like_title(first_line):
        title = first_line
        content_lines = lines[1:] if len(lines) > 1 else []
        content = '\n'.join(content_lines).strip()
        return title, content
    
    # No clear title, return all as content
    return "", text


def _looks_like_title(text: str) -> bool:
    """
    Determine if text looks like a title.
    
    Args:
        text: Text to evaluate
        
    Returns:
        True if text appears to be a title
    """
    if not text:
        return False
    
    # Short text that's all caps
    if text.isupper() and len(text) < 80:
        return True
    
    # Starts with chapter/section indicators
    if text.startswith(('Chapter', 'Section', 'Part', 'Article')):
        return True
    
    # Numbered sections
    if re.match(r'^\d+\.?\s+[A-Z]', text):
        return True
    
    # Short text without sentence ending
    if len(text) < 60 and not text.rstrip().endswith(('.', '!', '?')):
        return True
    
    return False


def _add_natural_content(content_placeholder, content_text: str) -> None:
    """
    Add content to placeholder in natural format.
    
    Args:
        content_placeholder: PowerPoint content placeholder
        content_text: Text content to add
    """
    text_frame = content_placeholder.text_frame
    text_frame.clear()
    
    # Split content into paragraphs
    paragraphs = content_text.split('\n\n')
    
    for i, paragraph_text in enumerate(paragraphs):
        paragraph_text = paragraph_text.strip()
        if not paragraph_text:
            continue
        
        if i == 0:
            # Use existing paragraph
            p = text_frame.paragraphs[0]
        else:
            # Add new paragraph
            p = text_frame.add_paragraph()
        
        p.text = paragraph_text
        p.font.name = 'Calibri'
        
        # Format based on content type
        if _looks_like_title(paragraph_text):
            p.font.size = Pt(18)
            p.font.bold = True
        elif paragraph_text.startswith(('•', '-', '*')) or re.match(r'^\d+\.', paragraph_text):
            p.font.size = Pt(14)
            p.level = 1  # Indent bullet points
        else:
            p.font.size = Pt(16)


def _add_manual_content_box(slide, content_text: str) -> None:
    """
    Add content as manual text box when no placeholder available.
    
    Args:
        slide: PowerPoint slide
        content_text: Text content to add
    """
    # Create text box for content
    left = Inches(0.5)
    top = Inches(1.5)
    width = Inches(9)
    height = Inches(5.5)
    
    textbox = slide.shapes.add_textbox(left, top, width, height)
    text_frame = textbox.text_frame
    text_frame.text = content_text
    text_frame.word_wrap = True
    text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
    
    # Format text
    for paragraph in text_frame.paragraphs:
        paragraph.font.name = 'Calibri'
        paragraph.font.size = Pt(16)


def calculate_optimal_slide_size(pdf_width: float, pdf_height: float) -> SlideConfig:
    """
    Calculate optimal slide dimensions based on PDF page size.
    Uses standard PowerPoint dimensions.
    
    Args:
        pdf_width: PDF page width in points
        pdf_height: PDF page height in points
        
    Returns:
        SlideConfig with standard dimensions
    """
    # Use standard widescreen format
    width_inches = 13.333  # Standard widescreen width
    height_inches = 7.5    # Standard widescreen height
    
    # Convert to EMU
    width_emu = int(width_inches * 914400)
    height_emu = int(height_inches * 914400)
    
    logger.info(f"Using standard slide size: {width_inches}\" x {height_inches}\"")
    
    return SlideConfig(width_emu, height_emu)


def create_empty_presentation(slide_count: int, slide_config: SlideConfig) -> bytes:
    """
    Create an empty presentation with the specified number of blank slides.
    
    Args:
        slide_count: Number of slides to create
        slide_config: Slide configuration
        
    Returns:
        PPTX file as bytes
    """
    try:
        prs = Presentation()
        
        # Add blank slides using title and content layout
        title_content_layout = prs.slide_layouts[1]
        for i in range(slide_count):
            prs.slides.add_slide(title_content_layout)
        
        # Save to bytes
        pptx_bytes = io.BytesIO()
        prs.save(pptx_bytes)
        pptx_bytes.seek(0)
        
        logger.info(f"Created empty presentation with {slide_count} slides")
        return pptx_bytes.getvalue()
        
    except Exception as e:
        logger.error(f"Failed to create empty presentation: {str(e)}")
        raise Exception(f"Failed to create empty presentation: {str(e)}")
